from pptx import Presentation
from langchain_core.documents import Document


def load_ppt(file_path: str):

    presentation = Presentation(file_path)

    documents = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1
    ):

        text = ""

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                text += shape.text + "\n"

        documents.append(
            Document(
                page_content=text,
                metadata={
                    "slide": slide_number,
                    "source": file_path
                }
            )
        )

    return documents